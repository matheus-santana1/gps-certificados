from pandas.tseries.offsets import CustomBusinessDay
from datetime import datetime
from pptx import Presentation
from openpyxl import load_workbook
from tkinter import messagebox, scrolledtext
from utils import converter_pasta_pdf
from docx import Document
import tkinter as tk
import pandas as pd
import holidays
import fitz 
import os
import sys
import psutil
import threading
import config

# --- FUNÇÕES DE AUTOMAÇÃO (MANTIDAS DO SEU ORIGINAL) ---

def substituir_pptx(modelo, dados_linha, pasta_pessoa):
    prs = Presentation(modelo)
    def processar_paragrafo(p):
        texto = "".join(run.text for run in p.runs)
        for k, v in dados_linha.items():
            texto = texto.replace(f"{{{{{k}}}}}", str(v))
        for run in p.runs: run.text = ""
        if p.runs: p.runs[0].text = texto
        else: p.add_run().text = texto
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs: processar_paragrafo(p)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for p in cell.text_frame.paragraphs: processar_paragrafo(p)
    nome = os.path.basename(modelo)
    prs.save(os.path.join(pasta_pessoa, nome))

def substituir_xlsx(modelo, dados_linha, pasta_pessoa):
    wb = load_workbook(modelo)
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                if isinstance(cell.value, str):
                    for k, v in dados_linha.items():
                        cell.value = cell.value.replace(f"{{{{{k}}}}}", str(v))
    nome = os.path.basename(modelo)
    wb.save(os.path.join(pasta_pessoa, nome))

def substituir_pdf(modelo, dados_linha, pasta_pessoa):
    if "DATA" in dados_linha and len(str(dados_linha["DATA"])) >= 5:
        dados_linha["DATA"] = str(dados_linha["DATA"])[:5]

    caminho_pdf_final = os.path.join(pasta_pessoa, "pdf")
    if not os.path.exists(caminho_pdf_final):
        os.makedirs(caminho_pdf_final)

    doc = fitz.open(modelo)
    for page in doc:
        widgets = page.widgets()
        for widget in widgets:
            if widget.field_name in dados_linha:
                widget.field_value = str(dados_linha[widget.field_name])
                widget.update() 
    nome_arquivo = os.path.basename(modelo)
    doc.save(os.path.join(caminho_pdf_final, nome_arquivo), incremental=False, encryption=fitz.PDF_ENCRYPT_KEEP)
    doc.close()

def substituir_docx(modelo, dados_linha, pasta_pessoa):
    doc = Document(modelo)
    def processar_paragrafo(paragrafo):
        if paragrafo.text:
            texto_atual = paragrafo.text
            alterado = False
            for k, v in dados_linha.items():
                placeholder = f"{{{{{k}}}}}"
                if placeholder in texto_atual:
                    texto_atual = texto_atual.replace(placeholder, str(v))
                    alterado = True
            if alterado:
                paragrafo.text = texto_atual
    for para in doc.paragraphs:
        processar_paragrafo(para)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    processar_paragrafo(para)
    nome = os.path.basename(modelo)
    caminho_final = os.path.join(pasta_pessoa, nome)
    doc.save(caminho_final)

def resource_path(relative_path):
    try:
        base_path = sys._MEIPASS
    except Exception:
        base_path = os.path.abspath(".")
    return os.path.join(base_path, relative_path)

def forcar_fechamento_seguro():
    for pid in config.pids_automacao:
        try:
            p = psutil.Process(pid)
            p.kill()
        except:
            pass
    config.pids_automacao.clear()

# --- LÓGICA PRINCIPAL E INTERFACE ---
def iniciar_automacao():
    def rodar():
        try:
            MODELOS = "modelos"
            SAIDA = "saida"
            DADOS = "base_de_dados.xlsx"
            
            # Verificações básicas
            if not os.path.exists(DADOS):
                app.log("Erro: Arquivo base_de_dados.xlsx não encontrado!")
                return
            if not os.path.exists(MODELOS):
                app.log("Erro: Pasta 'modelos' não encontrada!")
                return

            br_holidays = holidays.Brazil(state="MG")
            br = CustomBusinessDay(holidays=br_holidays)
            os.makedirs(SAIDA, exist_ok=True)
            
            dados = pd.read_excel(DADOS, skiprows=1)
            total = len(dados)

            for i, linha in dados.iterrows():
                dados_linha = linha.to_dict()
                
                # Tratamento de Datas
                data_base = pd.to_datetime(dados_linha["DATA DE ADMISSÃO"])
                dados_linha["DATA_1"] = data_base.strftime("%d/%m/%Y")
                for n in range(2, 6):
                    dados_linha[f"DATA_{n}"] = (data_base + br*(n-1)).strftime("%d/%m/%Y")
                
                dados_linha["DATA DE ADMISSÃO"] = data_base.strftime("%d/%m/%Y")
                dados_linha["PERÍODO DE ESPECÍFICO NA FUNÇÃO"] = f"{dados_linha['DATA_1']} à {dados_linha['DATA_5']}"
                
                if str(dados_linha.get("DATA")) == "(HOJE)":
                    dados_linha["DATA"] = datetime.now().strftime("%d/%m/%Y")

                nome_pessoa = str(dados_linha["NOME"])
                app.log(f"[{i+1}/{total}] Criando arquivos para: {nome_pessoa}")
                
                # Criar pasta
                pasta_nome = nome_pessoa.strip().replace("/", "-")
                pasta_pessoa = os.path.join(SAIDA, pasta_nome)
                os.makedirs(pasta_pessoa, exist_ok=True)

                # Processar modelos
                for arquivo in os.listdir(MODELOS):
                    caminho = os.path.join(MODELOS, arquivo)
                    if arquivo.endswith(".pptx"): substituir_pptx(caminho, dados_linha, pasta_pessoa)
                    elif arquivo.endswith(".xlsx"): substituir_xlsx(caminho, dados_linha, pasta_pessoa)
                    elif arquivo.endswith(".pdf"): substituir_pdf(caminho, dados_linha, pasta_pessoa)
                    elif arquivo.endswith(".docx") and arquivo[:-5] in dados_linha["ORDEM DE SERVIÇO"].split(","): substituir_docx(caminho, dados_linha, pasta_pessoa)
                    
            app.log("✅ Processo concluído com sucesso!")
            messagebox.showinfo("Sucesso", "Todos os arquivos foram gerados!")
        
        except Exception as e:
            app.log(f"❌ ERRO: {str(e)}")
            messagebox.showerror("Erro Fatal", f"Ocorreu um erro: {e}")
    threading.Thread(target=rodar, daemon=True).start()

def iniciar_conversao_pdf():
    DADOS = "base_de_dados.xlsx"
    config_excel = {}
    try:
        df = pd.read_excel(DADOS, skiprows=1)
        if "FRENTE_VERSO" in df.columns:
            valor_fv = df["FRENTE_VERSO"].dropna().iloc[0] if not df["FRENTE_VERSO"].dropna().empty else None
            if valor_fv and str(valor_fv).lower() != "nan":
                config_excel["frente_verso"] = [item.strip() for item in str(valor_fv).split(",")]

        for _, linha in df.iterrows():
            nome_pessoa = str(linha["NOME"]).strip()
            config_excel[nome_pessoa] = {}
            for coluna in df.columns:
                if str(coluna).lower().endswith(".xlsx"):
                    valor_celula = str(linha[coluna])
                    if valor_celula and valor_celula.lower() != "nan":
                        lista_abas = [aba.strip() for aba in valor_celula.split(",")]
                        config_excel[nome_pessoa][coluna] = lista_abas
    except Exception as e:
        app.log(f"⚠️ Aviso: Não foi possível ler as configurações ({e}). Prosseguindo sem filtros.")
    threading.Thread(target=lambda: converter_pasta_pdf("saida", app, config_excel), daemon=True).start()

class AppAutomação:
    def __init__(self, root):
        self.root = root
        self.root.title("GPS Certificados - Automação de Documentos")
        self.root.geometry("500x400")
        
        try:
            self.root.iconbitmap(resource_path("automation.ico"))
        except:
            pass

        # Frame para os botões ficarem lado a lado
        self.frame_topo = tk.Frame(root)
        self.frame_topo.pack(pady=20, padx=20, fill=tk.X)

        # Botão Iniciar
        self.btn_iniciar = tk.Button(
            self.frame_topo, 
            text="INICIAR", 
            command=iniciar_automacao, 
            bg="#4CAF50", 
            fg="white", 
            font=("Arial", 11, "bold"),
            pady=10
        )
        self.btn_iniciar.pack(side=tk.LEFT, fill=tk.X, expand=True)

        # Botão PDF
        self.btn_pdf = tk.Button(
            self.frame_topo, 
            text="📄", 
            command=iniciar_conversao_pdf, 
            bg="#2196F3", 
            fg="white", 
            font=("Arial", 14, "bold"),
            width=5,
            pady=7
        )
        self.btn_pdf.pack(side=tk.LEFT, padx=(10, 0))

        # Área de Log
        self.lbl_log = tk.Label(root, text="Status do Processamento:")
        self.lbl_log.pack(anchor="w", padx=20)
        self.txt_log = scrolledtext.ScrolledText(root, height=15, font=("Consolas", 9), state='disabled')
        self.txt_log.pack(pady=10, padx=20, fill=tk.BOTH, expand=True)
        self.root.protocol("WM_DELETE_WINDOW", self.ao_fechar)

    def log(self, mensagem):
        self.txt_log.config(state='normal')
        self.txt_log.insert(tk.END, mensagem + "\n")
        self.txt_log.see(tk.END)
        self.txt_log.config(state='disabled')

    def ao_fechar(self):
        if messagebox.askokcancel("Sair", "Deseja encerrar o programa? Isso interromperá as conversões em andamento."):
            forcar_fechamento_seguro()
            self.root.destroy()
            os._exit(0)

if __name__ == "__main__":
    root = tk.Tk()
    app = AppAutomação(root)
    root.mainloop()

